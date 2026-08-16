# -*- coding: utf-8 -*-
"""生成《小熊电器 AI 智能客服》方案 PPTX（12 页，16:9）。

配色：炭黑底 + 小熊橙，与 HTML 版一致，全部用原生可编辑形状。
运行：python build_pptx.py  → 输出 docs/方案PPT.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
import copy

# ---------- 配色 ----------
BG      = RGBColor(0x17, 0x11, 0x0C)   # 炭黑
CARD    = RGBColor(0x22, 0x1A, 0x13)   # 卡片深棕
CARD2   = RGBColor(0x2A, 0x20, 0x17)   # 卡片浅棕
ACCENT  = RGBColor(0xF4, 0x57, 0x1E)   # 小熊橙（检索/生成层）
SOFT    = RGBColor(0xFF, 0x8A, 0x4C)   # 暖橙（应用层）
GOLD    = RGBColor(0xE0, 0xA8, 0x4C)   # 金橙（对话层）
DARK    = RGBColor(0x1D, 0x16, 0x11)   # 深棕（面板底）
TEXT    = RGBColor(0xF6, 0xF0, 0xE6)   # 主文字
MUTED   = RGBColor(0xA9, 0x9C, 0x8B)   # 次要文字
LINE    = RGBColor(0x3A, 0x2E, 0x24)   # 分隔线

FONT = "Microsoft YaHei"
SW, SH = Inches(13.333), Inches(7.5)   # 16:9


def _ea_font(run, name):
    """设置东亚字体（PowerPoint 中文字体需写入 a:ea 元素，否则中文回退默认字体）。"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ---------- 工具函数 ----------
def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=CARD, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE, shadow=False):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=Pt(0), line_spacing=1.0):
    """runs: str 或 [(text, size, color, bold), ...] 列表"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, 18, TEXT, False)]
    for i, (t_, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = space_after
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = t_
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = FONT
        _ea_font(r, FONT)  # 东亚字体
    return tb


def shape_text(sp, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0):
    """在已有形状里填文字（多段落）。runs 可为 str 或段落列表。"""
    tf = sp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    if isinstance(runs, str):
        runs = [[(runs, 16, TEXT, False)]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for (t_, size, color, bold) in para:
            r = p.add_run()
            r.text = t_
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = FONT
            _ea_font(r, FONT)


def header(slide, title, sub=None):
    """页眉：橙色竖条 + 标题 + 可选副标题"""
    box(slide, Inches(0.6), Inches(0.5), Inches(0.09), Inches(0.62), fill=ACCENT)
    txt(slide, Inches(0.85), Inches(0.45), Inches(11.8), Inches(0.7),
        [(title, 30, TEXT, True)], anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        txt(slide, Inches(0.85), Inches(1.05), Inches(11.8), Inches(0.4),
            [(sub, 14, MUTED, False)])


def bullet(slide, l, t, w, items, size=16, gap=Pt(8), color=TEXT, marker="· "):
    """items: [(主文本, 次文本), ...] 或 [str, ...]"""
    paras = []
    for it in items:
        if isinstance(it, tuple):
            main, sub = it
            paras.append([(marker + main, size, color, True)])
            if sub:
                paras.append([("    " + sub, size - 3, MUTED, False)])
        else:
            paras.append([(marker + it, size, color, False)])
    tb = slide.shapes.add_textbox(l, t, w, Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        p.line_spacing = 1.15
        for (t_, sz, clr, b) in para:
            r = p.add_run()
            r.text = t_
            r.font.size = Pt(sz)
            r.font.color.rgb = clr
            r.font.bold = b
            r.font.name = FONT
            _ea_font(r, FONT)
    return tb


# ============================================================
# 第 1 页 · 封面
# ============================================================
s = add_slide(); set_bg(s)
box(s, Inches(0), Inches(0), SW, SH, fill=BG)
box(s, Inches(0), Inches(7.02), SW, Inches(0.48), fill=ACCENT)  # 底部色带
txt(s, Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.6),
    [("小熊电器 AI 智能客服", 54, TEXT, True)], align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(3.35), Inches(11.3), Inches(0.6),
    [("让 AI 从「人工辅助」走向「全自主接待」", 24, SOFT, False)], align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(4.35), Inches(11.3), Inches(0.6),
    [("基于大模型 RAG，融合产品知识库与多平台规则", 17, MUTED, False)], align=PP_ALIGN.CENTER)
# 装饰点
for i, c in enumerate([ACCENT, SOFT, MUTED]):
    box(s, Inches(6.07 + i * 0.42), Inches(5.35), Inches(0.22), Inches(0.22),
        fill=c, shape=MSO_SHAPE.OVAL)

# ============================================================
# 第 2 页 · 痛点
# ============================================================
s = add_slide(); set_bg(s)
header(s, "为什么要做？", "千牛人工分岗 + 转接流转，五大痛点")
pains = [
    ("分流不精准、转接多", "售前/售中人工分岗，分流不准就转接，用户被转来转去、转化流失"),
    ("接待受人力限制", "淘系/京东/抖音多平台并线，高峰响应慢，体验分与转化率承压"),
    ("话术标准化难", "售前/售中话术参差，新人培训 2-4 周，知识更新滞后"),
    ("质量靠事后记录", "赤兔名品只记接待量/满意度/转化率，无法实时提升话术"),
    ("多平台规则适配难", "京东/天猫/抖音规则各异，人工易记混出错"),
]
for i, (title, desc) in enumerate(pains):
    if i == 4:  # 第 5 项居中
        l, t, w, h = Inches(3.9), Inches(5.35), Inches(5.5), Inches(1.65)
    else:
        r, c = divmod(i, 2)
        l, t, w, h = Inches(0.7 + c * 6.2), Inches(1.6 + r * 1.9), Inches(5.8), Inches(1.75)
    box(s, l, t, w, h, fill=CARD, line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, l, t, Inches(0.09), h, fill=ACCENT)
    txt(s, l + Inches(0.3), t + Inches(0.2), w - Inches(0.5), Inches(0.5),
        [(f"0{i+1}  {title}", 19, TEXT, True)])
    txt(s, l + Inches(0.3), t + Inches(0.82), w - Inches(0.5), Inches(0.85),
        [(desc, 14, MUTED, False)], line_spacing=1.1)

# ============================================================
# 第 3 页 · 我们的方案
# ============================================================
s = add_slide(); set_bg(s)
header(s, "我们的方案", "一句话讲清核心思路")
box(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.6), fill=CARD, line=SOFT, line_w=Pt(1.5),
    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(box(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.6), fill=None),
    [[("一个 AI，全岗位自主接待", 30, TEXT, True)],
     [("替代「人工分岗 + 分流 + 转接」", 22, ACCENT, True)]])
txt(s, Inches(1.0), Inches(4.6), Inches(11.3), Inches(0.6),
    [("进线即自动分流（售前 / 售中 / 售后）→ RAG 检索 + 大模型生成 → 秒级回复、无需转接", 16, MUTED, False)],
    align=PP_ALIGN.CENTER)

# ============================================================
# 第 4 页 · 对标千牛工作流：分流与转接
# ============================================================
s = add_slide(); set_bg(s)
header(s, "对标千牛工作流", "从「人工分岗 + 转接」到「AI 全岗位自主接待」")

RED = RGBColor(0xD1, 0x4A, 0x3C)


def f_node(x, y, w, h, big, sub, color, hl=False):
    """流程节点：左侧色条 + 大字 + 副标题。"""
    box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=CARD,
        line=(color if hl else LINE), line_w=Pt(1.5), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, Inches(x), Inches(y), Inches(0.09), Inches(h), fill=color)
    shape_text(box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=None),
               [[(big, 14, (color if hl else TEXT), True)], [(sub, 10, MUTED, False)]])


def f_arrow(x, y):
    box(s, Inches(x), Inches(y), Inches(0.4), Inches(0.3), fill=ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)


# ---- 现状：千牛人工 ----
txt(s, Inches(0.7), Inches(1.5), Inches(4.5), Inches(0.5),
    [("现状 · 千牛人工", 16, MUTED, True), ("　分岗 + 转接", 13, MUTED, False)])
old_flow = [
    ("客户进线", "淘系/京东/抖音", MUTED, False),
    ("系统分流", "人工/规则分流", MUTED, False),
    ("售前 / 售中", "分岗位接待", MUTED, False),
    ("转接 / 流失", "分流不精准", RED, True),
]
for i, (big, sub, c, hl) in enumerate(old_flow):
    f_node(0.7 + i * 3.1, 2.1, 2.6, 1.05, big, sub, c, hl)
    if i < 3:
        f_arrow(3.4 + i * 3.1, 2.5)
txt(s, Inches(0.7), Inches(3.28), Inches(11.9), Inches(0.4),
    [("分流不精准 → 转接 → 用户被转来转去、转化流失；质量只能事后靠赤兔名品记录", 12.5, RED, False)])

# ---- 方案：AI 自主接待 ----
txt(s, Inches(0.7), Inches(3.95), Inches(4.5), Inches(0.5),
    [("方案 · AI 自主接待", 16, SOFT, True), ("　进线即自动分流", 13, MUTED, False)])
new_flow = [
    ("客户进线", "多平台统一网关", SOFT, False),
    ("AI 自动分流", "售前/售中/售后", SOFT, False),
    ("RAG 检索 + 生成", "知识库 + 大模型", ACCENT, False),
    ("秒级回复", "无需转接", ACCENT, True),
]
for i, (big, sub, c, hl) in enumerate(new_flow):
    f_node(0.7 + i * 3.1, 4.55, 2.6, 1.05, big, sub, c, hl)
    if i < 3:
        f_arrow(3.4 + i * 3.1, 4.95)
txt(s, Inches(0.7), Inches(5.72), Inches(11.9), Inches(0.4),
    [("AI 进线即自动分流到售前 / 售中 / 售后，直接解决咨询，全程无需转接", 12.5, ACCENT, False)])

# ---- 底部：赤兔名品三指标 ----
box(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.85), fill=DARK, line=LINE, line_w=Pt(1),
    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(box(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.85), fill=None),
    [[("赤兔名品三指标", 13, GOLD, True),
      ("　接待量 · 满意度 · 转化率", 12, MUTED, False),
      ("　→　现状仅事后记录，", 11, MUTED, False),
      ("AI 实时提升", 11, SOFT, True)]])

# ============================================================
# 第 5 页 · 核心能力
# ============================================================
s = add_slide(); set_bg(s)
header(s, "核心能力", "对应赛题的核心要求")
cards = [
    ("① 话术生成", ["RAG 检索产品知识 + 平台规则", "实时生成贴合问题的回复话术", "三要素：卖点(FAB) + 对比 + 促单"]),
    ("② 岗位自动分流", ["进线即意图识别", "自动路由 售前 / 售中 / 售后", "无需人工分岗与转接"]),
    ("③ 智能接待", ["多轮对话上下文理解", "主动引导、推荐、挖需、促单", "跑通「推荐→讲解→挖需→促单」"]),
]
for i, (t, lines) in enumerate(cards):
    l = Inches(0.7 + i * 4.1)
    w = Inches(3.75)
    box(s, l, Inches(1.9), w, Inches(3.4), fill=CARD, line=LINE, line_w=Pt(1),
        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, l, Inches(1.9), w, Inches(0.75), fill=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, l + Inches(0.3), Inches(2.05), w - Inches(0.6), Inches(0.5), [(t, 20, SOFT, True)])
    bullet(s, l + Inches(0.3), Inches(2.75), w - Inches(0.6),
           [(ln, "") for ln in lines], size=14, color=TEXT, gap=Pt(10))

# ============================================================
# 第 6 页 · 技术架构（原生可编辑形状，五层颜色编码）
# ============================================================
s = add_slide(); set_bg(s)
header(s, "技术架构", "数据层 → 检索层 → 生成层 → 对话层 → 应用层")

def pipe_box(x, y, w, h, title, sub, layer, lc, fill=CARD):
    """主流程方框：左侧层级色条 + 标题/副标题 + 右上角层级标签。"""
    box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=fill, line=lc, line_w=Pt(1.5),
        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, Inches(x), Inches(y), Inches(0.09), Inches(h), fill=lc)   # 左侧色条
    runs = [[(title, 16, TEXT, True)]]
    if sub:
        runs.append([(sub, 10.5, MUTED, False)])
    shape_text(box(s, Inches(x + 0.2), Inches(y), Inches(w - 0.2), Inches(h), fill=None), runs)
    tag = box(s, Inches(x + w - 0.82), Inches(y + 0.08), Inches(0.72), Inches(0.28),
              fill=lc, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(tag, [[(layer, 9.5, BG, True)]])


# 主流程竖排（居中）
CX, W = 4.55, 4.2
pipe_box(CX, 1.55, W, 0.6,  "买家咨询",     None,                     "应用层", SOFT,  fill=CARD2)
pipe_box(CX, 2.45, W, 0.72, "意图识别",     "推荐·对比·参数·优惠·物流·售后", "对话层", GOLD)
pipe_box(CX, 3.47, W, 0.72, "检索 · BM25",  "中文分词 · Top-K 召回",    "检索层", ACCENT)
pipe_box(CX, 4.49, W, 0.72, "生成 · DeepSeek", "卖点(FAB)+对比+促单",   "生成层", ACCENT)
pipe_box(CX, 5.51, W, 0.6,  "AI 回复话术",  None,                     "应用层", SOFT,  fill=CARD2)

# 向下流程箭头
for y in [2.17, 3.19, 4.21, 5.23]:
    box(s, Inches(CX + W/2 - 0.14), Inches(y), Inches(0.28), Inches(0.24),
        fill=ACCENT, shape=MSO_SHAPE.DOWN_ARROW)

# 左侧：多轮对话管理（指向生成层）
lx, ly, lw, lh = 0.45, 4.1, 3.6, 1.5
box(s, Inches(lx), Inches(ly), Inches(lw), Inches(lh), fill=DARK, line=MUTED, line_w=Pt(1),
    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
box(s, Inches(lx), Inches(ly), Inches(0.09), Inches(lh), fill=GOLD)
shape_text(box(s, Inches(lx + 0.2), Inches(ly), Inches(lw - 0.2), Inches(lh), fill=None),
    [[("多轮对话管理", 14, GOLD, True)],
     [("会话历史 · 漏斗阶段跟踪", 10.5, MUTED, False)],
     [("推荐→讲解→挖需→促单", 10.5, MUTED, False)]])
box(s, Inches(lx + lw + 0.02), Inches(ly + lh/2 - 0.13), Inches(CX - lx - lw - 0.06), Inches(0.26),
    fill=GOLD, shape=MSO_SHAPE.RIGHT_ARROW)

# 右侧：数据层（指向检索层）
rx, ry, rw, rh = 9.3, 2.9, 3.5, 1.85
box(s, Inches(rx), Inches(ry), Inches(rw), Inches(rh), fill=DARK, line=MUTED, line_w=Pt(1),
    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(rx + 0.2), Inches(ry + 0.08), Inches(rw - 0.4), Inches(0.32),
    [("数据层", 14, SOFT, True)], align=PP_ALIGN.CENTER)
box(s, Inches(rx + 0.16), Inches(ry + 0.42), Inches(rw - 0.32), Inches(0.62),
    fill=CARD2, line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(box(s, Inches(rx + 0.16), Inches(ry + 0.42), Inches(rw - 0.32), Inches(0.62), fill=None),
    [[("产品知识库", 12, TEXT, True)], [("卖点 / 参数 / FAQ", 9.5, MUTED, False)]])
box(s, Inches(rx + 0.16), Inches(ry + 1.1), Inches(rw - 0.32), Inches(0.62),
    fill=CARD2, line=LINE, line_w=Pt(1), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(box(s, Inches(rx + 0.16), Inches(ry + 1.1), Inches(rw - 0.32), Inches(0.62), fill=None),
    [[("平台规则库", 12, TEXT, True)], [("优惠 / 售后 / 风格", 9.5, MUTED, False)]])
box(s, Inches(rx - 0.38), Inches(ry + 0.72 - 0.13), Inches(0.38), Inches(0.26),
    fill=SOFT, shape=MSO_SHAPE.LEFT_ARROW)

# 底部说明
box(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.72), fill=DARK, line=LINE, line_w=Pt(1),
    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
shape_text(box(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.72), fill=None),
    [[("应用层", 11, SOFT, True), ("  Streamlit 可演示 Demo（话术生成 / 智能接待 / 批量对比）　", 11, MUTED, False),
      ("数据层", 11, SOFT, True), ("  改 JSON 即更新，零训练成本", 11, MUTED, False)]])

# ============================================================
# 第 7 页 · Demo A
# ============================================================
s = add_slide(); set_bg(s)
header(s, "现场演示 · Demo A", "话术生成")
box(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(0.9), fill=CARD, line=LINE, line_w=Pt(1),
    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.0), Inches(2.02), Inches(11.3), Inches(0.66),
    [("输入：", 15, SOFT, True), ("“这款按摩器和 SKG 比有什么优势？”", 16, TEXT, False)],
    anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(0.7), Inches(3.05), Inches(11.9), Inches(2.6), fill=CARD, line=ACCENT, line_w=Pt(1.2),
    shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.0), Inches(3.2), Inches(11.3), Inches(0.5),
    [("输出话术（京东风格）", 15, ACCENT, True)])
bullet(s, Inches(1.0), Inches(3.75), Inches(11.3), [
    ("卖点", "「3D 仿人手揉捏 + 恒温热敷，缓解颈部疲劳」"),
    ("对比", "「相比 SKG 同档位便宜约 30%，配置不缩水」"),
    ("促单", "「京东自营当日发，满 199 减 30，下单赠收纳袋」"),
], size=15, gap=Pt(6))
txt(s, Inches(0.7), Inches(5.85), Inches(11.9), Inches(0.5),
    [("关键：展示「检索命中 → 上下文增强 → 生成」的可解释过程", 14, MUTED, False)])

# ============================================================
# 第 8 页 · Demo B
# ============================================================
s = add_slide(); set_bg(s)
header(s, "现场演示 · Demo B", "多轮智能接待（3 轮自主推进）")
steps = [
    ("买家", "「想给爸妈买个养生壶」", "AI 推荐 + 挖需", ACCENT),
    ("买家", "「能预约吗？」", "AI 讲解参数", SOFT),
    ("买家", "「有优惠吗？」", "AI 促单", ACCENT),
]
for i, (who, q, a, c) in enumerate(steps):
    t = Inches(1.5 + i * 1.75)
    box(s, Inches(0.7), t, Inches(11.9), Inches(1.45), fill=CARD, line=LINE, line_w=Pt(1),
        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, Inches(0.7), t, Inches(0.7), Inches(1.45), fill=c)
    txt(s, Inches(1.55), t + Inches(0.15), Inches(2.6), Inches(1.1),
        [(f"第 {i+1} 轮", 13, MUTED, False), ("", 6, MUTED, False), (who, 17, TEXT, True)])
    txt(s, Inches(3.2), t + Inches(0.18), Inches(5.0), Inches(1.1),
        [(q, 14, TEXT, False)])
    txt(s, Inches(8.3), t + Inches(0.18), Inches(4.0), Inches(1.1),
        [(a, 14, c, True)])
txt(s, Inches(0.7), Inches(6.9), Inches(11.9), Inches(0.5),
    [("自主完成「推荐 → 讲解 → 挖需 → 促单」全链路", 14, MUTED, False)])

# ============================================================
# 第 9 页 · Demo C
# ============================================================
s = add_slide(); set_bg(s)
header(s, "现场演示 · Demo C", "对齐赤兔名品三指标：接待量 / 满意度 / 转化率")
# 表头
box(s, Inches(0.7), Inches(1.7), Inches(3.3), Inches(0.6), fill=CARD2, line=LINE)
shape_text(box(s, Inches(0.7), Inches(1.7), Inches(3.3), Inches(0.6), fill=None), [[("维度", 15, MUTED, True)]])
box(s, Inches(4.0), Inches(1.7), Inches(4.4), Inches(0.6), fill=ACCENT, line=None)
shape_text(box(s, Inches(4.0), Inches(1.7), Inches(4.4), Inches(0.6), fill=None), [[("AI 话术", 15, TEXT, True)]])
box(s, Inches(8.4), Inches(1.7), Inches(4.2), Inches(0.6), fill=CARD2, line=LINE)
shape_text(box(s, Inches(8.4), Inches(1.7), Inches(4.2), Inches(0.6), fill=None), [[("人工话术", 15, MUTED, True)]])
rows = [
    ("接待量", "7×24 不设上限", "6-10 人并发，高峰排队"),
    ("满意度", "话术标准化、稳定", "因人而异、参差不齐"),
    ("转化率", "秒回 + 促单 + 无转接流失", "转接流失、响应慢"),
    ("响应速度", "秒级回复", "高峰排队 / 下班离线"),
    ("成本", "边际成本趋近 0", "随业务线性增长"),
]
for r, (dim, ai, human) in enumerate(rows):
    y = Inches(2.35 + r * 0.78)
    box(s, Inches(0.7), y, Inches(3.3), Inches(0.72), fill=CARD, line=LINE)
    shape_text(box(s, Inches(0.7), y, Inches(3.3), Inches(0.72), fill=None), [[(dim, 14, TEXT, True)]])
    box(s, Inches(4.0), y, Inches(4.4), Inches(0.72), fill=RGBColor(0x2A, 0x1F, 0x14), line=None)
    shape_text(box(s, Inches(4.0), y, Inches(4.4), Inches(0.72), fill=None), [[(ai, 13, SOFT, False)]])
    box(s, Inches(8.4), y, Inches(4.2), Inches(0.72), fill=CARD, line=LINE)
    shape_text(box(s, Inches(8.4), y, Inches(4.2), Inches(0.72), fill=None), [[(human, 13, MUTED, False)]])

# ============================================================
# 第 10 页 · 多平台差异化 + 统一接入网关
# ============================================================
s = add_slide(); set_bg(s)
header(s, "多平台适配", "一个模型 + 规则注入，自动切换四平台话术风格")
plats = [
    ("京东", "专业可靠", "自营正品 · 次日达", ACCENT),
    ("天猫", "亲切质感", "旗舰店 · 88VIP", SOFT),
    ("抖音", "热情口语", "限时福利 · 抢购", SOFT),
    ("拼多多", "接地气", "百亿补贴 · 性价比", ACCENT),
]
for i, (name, style, desc, c) in enumerate(plats):
    l = Inches(0.7 + i * 3.05)
    t = Inches(1.55)
    box(s, l, t, Inches(2.9), Inches(1.55), fill=CARD, line=LINE, line_w=Pt(1),
        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, l, t, Inches(2.9), Inches(0.55), fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, l + Inches(0.2), t + Inches(0.07), Inches(2.5), Inches(0.4), [(name, 18, TEXT, True)])
    txt(s, l + Inches(0.2), t + Inches(0.66), Inches(2.5), Inches(0.35), [(style, 14, SOFT, True)])
    txt(s, l + Inches(0.2), t + Inches(1.05), Inches(2.5), Inches(0.42), [(desc, 11.5, MUTED, False)])


def g_node(x, y, w, h, big, sub, hl=False):
    """统一网关流程节点。"""
    box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=CARD,
        line=(ACCENT if hl else LINE), line_w=Pt(1.25), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    shape_text(box(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=None),
               [[(big, 13, (SOFT if hl else TEXT), True)], [(sub, 10.5, MUTED, False)]])


nodes = [
    (0.7,  "四平台",     "买家消息",          False),
    (3.8,  "统一网关",   "/webhook/{平台}",   True),
    (6.9,  "平台适配器", "Adapter",          True),
    (10.0, "RAG Engine", "话术生成",         True),
]
for (x, big, sub, hl) in nodes:
    g_node(x, 3.55, 2.6, 1.2, big, sub, hl)
for x in [3.3, 6.4, 9.5]:
    box(s, Inches(x), Inches(4.0), Inches(0.5), Inches(0.3), fill=ACCENT, shape=MSO_SHAPE.RIGHT_ARROW)

txt(s, Inches(0.7), Inches(5.15), Inches(11.9), Inches(0.4),
    [("话术经同一网关回发至对应平台　", 13, MUTED, False),
     ("平台差异只隔离在 Adapter 层", 13, SOFT, True)])
txt(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.5),
    [("规则独立更新，零训练成本 —— 直接回应「多平台规则适配难」痛点", 14, ACCENT, True)])

# ============================================================
# 第 11 页 · 价值与落地
# ============================================================
s = add_slide(); set_bg(s)
header(s, "价值与落地", "对赤兔名品三大指标的提升")
vals = [
    ("接待量 ↑", "AI 7×24 不设上限，人工转异常兜底"),
    ("满意度 ↑", "话术标准化，质量稳定不随人波动"),
    ("转化率 ↑", "秒回 + 促单 + 无转接流失"),
    ("易落地", "改 JSON 即更新，零训练成本"),
]
for i, (t, d) in enumerate(vals):
    r, c = divmod(i, 2)
    l = Inches(0.7 + c * 6.2)
    tt = Inches(1.9 + r * 2.2)
    box(s, l, tt, Inches(5.8), Inches(1.9), fill=CARD, line=LINE, line_w=Pt(1),
        shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, l, tt, Inches(0.09), Inches(1.9), fill=ACCENT)
    txt(s, l + Inches(0.35), tt + Inches(0.25), Inches(5.2), Inches(0.5), [(t, 22, SOFT, True)])
    txt(s, l + Inches(0.35), tt + Inches(0.9), Inches(5.2), Inches(0.8), [(d, 14, MUTED, False)],
        line_spacing=1.1)

# ============================================================
# 第 12 页 · 总结
# ============================================================
s = add_slide(); set_bg(s)
box(s, Inches(0), Inches(0), SW, SH, fill=BG)
box(s, Inches(0), Inches(0), SW, Inches(0.4), fill=ACCENT)
txt(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.4),
    [("让 AI 客服从「辅助工具」", 38, TEXT, True)],
    align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(3.3), Inches(11.3), Inches(1.4),
    [("升级为「全自主接待单元」", 38, ACCENT, True)],
    align=PP_ALIGN.CENTER)
txt(s, Inches(1.0), Inches(4.9), Inches(11.3), Inches(0.6),
    [("跑通「买家提问 → RAG 检索 → 话术生成 → 自主回复」无人化链路", 16, MUTED, False)],
    align=PP_ALIGN.CENTER)

# ---------- 保存 ----------
out = r"C:\Users\夏\Desktop\bear-rag-cs\docs\方案PPT.pptx"
prs.save(out)
print("已生成：", out)
print("页数：", len(prs.slides._sldIdLst))
